#!/usr/bin/env python3
"""Generate Belfiore Cheese pitch deck (.pptx) -- McKinsey SCR narrative.

12 slides: Cover -> Exec Summary -> Market (S) -> Company (S) ->
Succession Gap (C) -> Nicholas Question (C) -> Growth Channels (R) ->
Business Economics (R) -> Personal Decision (R) -> Entry Strategy (R) ->
Risks (R) -> Next Steps

This is a PRIVATE analysis tool -- never shown to Farr or the Hariri family.
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from lxml import etree

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_DIR = os.path.join(REPO, "output", "belfiore-cheese")
OUT_FILE = os.path.join(OUT_DIR, "030326-pitch-deck.pptx")
ASSUMPTIONS_FILE = os.path.join(OUT_DIR, "assumptions.json")

# ── Colors ───────────────────────────────────────────────────────────────────
MCK_DARK_BLUE = RGBColor(0x05, 0x1C, 0x2C)
MCK_BLUE      = RGBColor(0x22, 0x51, 0xFF)
CORP_BLUE     = RGBColor(0x00, 0x5E, 0xB8)
NEAR_BLACK    = RGBColor(0x22, 0x22, 0x22)
DARK_GRAY     = RGBColor(0x58, 0x59, 0x5B)
MEDIUM_GRAY   = RGBColor(0xA2, 0xAA, 0xAD)
LIGHT_GRAY    = RGBColor(0xD9, 0xD9, 0xD9)
BG_GRAY       = RGBColor(0xF5, 0xF5, 0xF5)
GOLD          = RGBColor(0xF3, 0xC1, 0x3A)
ALERT_RED     = RGBColor(0xE3, 0x1B, 0x23)
GROWTH_GREEN  = RGBColor(0x00, 0x7A, 0x33)
CHEESE_GOLD   = RGBColor(0xF5, 0xDE, 0xB3)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
FOOTER_GRAY   = RGBColor(0x80, 0x80, 0x80)

# ── Dimensions ───────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TITLE_TOP    = Inches(0);     TITLE_H = Inches(1.1)
SUBTITLE_TOP = Inches(1.15);  SUBTITLE_H = Inches(0.3)
BODY_TOP     = Inches(1.55);  BODY_H = Inches(4.7)
SOURCE_TOP   = Inches(6.35);  SOURCE_H = Inches(0.3)
FOOTER_TOP   = Inches(6.85);  FOOTER_H = Inches(0.25)
CONTENT_L    = Inches(0.5);   CONTENT_W = Inches(12.333)


# ── Reusable components ─────────────────────────────────────────────────────

def mck_slide(prs, action_title, source_text, slide_num, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0), top=TITLE_TOP, width=SLIDE_W, height=TITLE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = MCK_DARK_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.25)
    tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = action_title
    p.alignment = PP_ALIGN.LEFT
    p.font.name = "Georgia"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        box = slide.shapes.add_textbox(CONTENT_L, SUBTITLE_TOP, CONTENT_W, SUBTITLE_H)
        stf = box.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Arial"
        sp.font.size = Pt(11)
        sp.font.color.rgb = DARK_GRAY

    src = slide.shapes.add_textbox(CONTENT_L, SOURCE_TOP, Inches(10), SOURCE_H)
    stf = src.text_frame
    sp = stf.paragraphs[0]
    sp.text = source_text
    sp.font.name = "Arial"
    sp.font.size = Pt(7)
    sp.font.italic = True
    sp.font.color.rgb = FOOTER_GRAY

    pn = slide.shapes.add_textbox(Inches(12.4), FOOTER_TOP, Inches(0.6), FOOTER_H)
    ptf = pn.text_frame
    pp = ptf.paragraphs[0]
    pp.text = str(slide_num)
    pp.alignment = PP_ALIGN.RIGHT
    pp.font.name = "Arial"
    pp.font.size = Pt(8)
    pp.font.color.rgb = FOOTER_GRAY

    cl = slide.shapes.add_textbox(CONTENT_L, FOOTER_TOP, Inches(4), FOOTER_H)
    ctf = cl.text_frame
    cp = ctf.paragraphs[0]
    cp.text = "PRIVATE & CONFIDENTIAL -- NOT FOR EXTERNAL USE"
    cp.font.name = "Arial"
    cp.font.size = Pt(7)
    cp.font.color.rgb = ALERT_RED

    return slide


def add_body_bullets(slide, bullets, top=None, left=None, width=None, height=None,
                     font_size=Pt(12)):
    t = top or BODY_TOP
    l = left or CONTENT_L
    w = width or CONTENT_W
    h = height or BODY_H

    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, is_bold = item
        else:
            text, is_bold = item, False

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if text.startswith(">>"):
            text = text[2:].strip()
            is_bold = True
            p.space_before = Pt(12)

        p.text = text
        p.font.name = "Arial"
        p.font.size = font_size
        p.font.color.rgb = NEAR_BLACK
        p.font.bold = is_bold
        p.space_after = Pt(4)
        p.line_spacing = Pt(18) if font_size <= Pt(12) else Pt(22)

    return box


def set_cell_border(cell, side, color_hex="D9D9D9", width_emu=6350):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    border_el = tcPr.find(qn(f"a:{side}"))
    if border_el is None:
        border_el = etree.SubElement(tcPr, qn(f"a:{side}"))
    border_el.set("w", str(width_emu))
    border_el.set("cap", "flat")
    border_el.set("cmpd", "sng")
    sf = etree.SubElement(border_el, qn("a:solidFill"))
    srgb = etree.SubElement(sf, qn("a:srgbClr"))
    srgb.set("val", color_hex)


def add_mck_table(slide, data, top=None, left=None, width=None,
                  col_widths=None, font_size=Pt(10)):
    t = top or BODY_TOP
    l = left or CONTENT_L
    w = width or CONTENT_W

    rows = len(data)
    cols = len(data[0])
    row_h = Inches(0.38)

    tbl_shape = slide.shapes.add_table(rows, cols, l, t, w, Emu(int(row_h) * rows))
    table = tbl_shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw

    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = font_size

                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = MCK_DARK_BLUE
                else:
                    paragraph.font.color.rgb = NEAR_BLACK
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else BG_GRAY

            for side in ["lnL", "lnR", "lnT"]:
                set_cell_border(cell, side, color_hex="FFFFFF", width_emu=0)
            set_cell_border(cell, "lnB", color_hex="D9D9D9", width_emu=6350)

    return tbl_shape


def add_gold_callout(slide, text, top=None, left=None, width=None):
    t = top or Inches(5.8)
    l = left or CONTENT_L
    w = width or CONTENT_W

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = MCK_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    return shape


def style_chart_mck(chart, show_value_axis=False, show_legend=False, gap_width=100):
    chart.has_legend = show_legend
    if show_legend and chart.has_legend:
        chart.legend.font.name = "Arial"
        chart.legend.font.size = Pt(9)
        chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.gap_width = gap_width

    va = chart.value_axis
    va.visible = show_value_axis
    va.has_major_gridlines = False
    va.has_minor_gridlines = False

    ca = chart.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.name = "Arial"
    ca.tick_labels.font.size = Pt(9)
    ca.tick_labels.font.color.rgb = DARK_GRAY
    ca.format.line.color.rgb = LIGHT_GRAY
    ca.format.line.width = Pt(0.5)


# ═════════════════════════════════════════════════════════════════════════════
# 12 SLIDES
# ═════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = MCK_DARK_BLUE

    box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.333), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Belfiore Cheese Company"
    p.alignment = PP_ALIGN.LEFT
    p.font.name = "Georgia"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box2 = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11.333), Inches(0.8))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Operating Entry Analysis -- Private Decision Document"
    p2.alignment = PP_ALIGN.LEFT
    p2.font.name = "Georgia"
    p2.font.size = Pt(24)
    p2.font.color.rgb = GOLD

    box3 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(0.6))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "March 2026  |  PRIVATE & CONFIDENTIAL -- NOT FOR EXTERNAL USE"
    p3.alignment = PP_ALIGN.LEFT
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = MEDIUM_GRAY

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.5), Inches(4.333), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def slide_02_exec_summary(prs):
    slide = mck_slide(prs,
        "Belfiore presents a real but narrow operating entry opportunity -- the most "
        "likely outcome is a valuable bridge role, not ownership",
        "Source: 9-agent research synthesis; family intelligence; case study analysis",
        slide_num=2)

    bullets = [
        ("THESIS", True),
        "Belfiore Cheese ($3-5M revenue, Berkeley) has an emerging succession gap. "
        "Farr Hariri (73-75, 35+ years) hasn't formalized a plan, but his son Nicholas "
        "is showing interest from SF. An outside operator with McKinsey + CoS credentials "
        "could enter through the family connection and position for a gradual transition.",
        "",
        ("PROBABILITY-WEIGHTED OUTCOMES", True),
        "\u2022  Best case (20%): Written buyout over 5-7 years. Nick owns Belfiore by 35. SDE $600K+.",
        "\u2022  Base case (40%): COO for 3-5 years, $130-160K comp, no equity. Nicholas enters; Nick exits.",
        "\u2022  Downside (25%): Never gets past relationship phase. Good connection, no career disruption.",
        "\u2022  Worst case (15%): 2-3 years at below-market, Nicholas takes over, Nick exits with nothing.",
        "",
        ("CRITICAL INSIGHT FROM CASE STUDIES", True),
        "\u2022  Zero documented cases of outside operators earning equity through sweat equity at artisan cheese scale.",
        "\u2022  Every successful ownership transition was direct acquisition, investor buyout, or corporate sale.",
        "\u2022  Verbal equity promises = zero value. Written agreement within 6 months or exit.",
    ]
    add_body_bullets(slide, bullets, font_size=Pt(11))


def slide_03_market(prs):
    slide = mck_slide(prs,
        "US specialty cheese is a $6.67B market growing 5.6% annually -- artisan "
        "producers are outperforming mass-market and taking share",
        "Source: USDA; Specialty Food Association; Grand View Research; trade publications",
        slide_num=3,
        subtitle="SITUATION -- Industry landscape")

    bullets = [
        "\u2022  US specialty cheese market: $6.67B (2025), projected $9.2B by 2030 at 5.6% CAGR",
        "\u2022  Artisan segment growing faster than commodity at 7-9% -- driven by clean-label, local sourcing, premiumization",
        "\u2022  EU tariffs (2025-2026) creating domestic opportunity: Italian imports face 25% tariffs, "
        "benefiting domestic mozzarella and burrata producers like Belfiore",
        "",
        "\u2022  Bay Area specifics: highest per-capita specialty cheese consumption in US; Whole Foods, "
        "Bi-Rite, Berkeley Bowl driving demand; consumer willingness to pay 30-50% premium for local artisan",
        "",
        ("\u2022  M&A trend: European dairy giants acquiring CA artisan brands at 5-7x EBITDA "
         "(Emmi/Cowgirl Creamery 2016, Savencia/Rogue Creamery 2018) -- validates exit multiples"),
        "",
        "\u2022  Headwinds: input cost inflation (milk, labor), Berkeley property pressure, "
        "distribution complexity at small scale (<$10M revenue)",
    ]
    add_body_bullets(slide, bullets, font_size=Pt(11))


def slide_04_company(prs):
    slide = mck_slide(prs,
        "Belfiore has operated profitably for 37 years with 25 years of national awards -- "
        "a durable craft business with a founder approaching the end of his tenure",
        "Source: Company research dossier; ACS competition records; CCOF directory; Gemini Deep Research",
        slide_num=4,
        subtitle="SITUATION -- Belfiore Cheese Company profile")

    data = [
        ["Attribute", "Detail"],
        ["Founded", "1987 (San Rafael); current ownership since 1989"],
        ["Revenue", "$3-5M estimated (ZoomInfo: $1-10M range)"],
        ["Employees", "~10-14 (flat, no growth hiring)"],
        ["Products", "13+ SKUs: Fior di Latte, Burrata, Burrata di Bufala,\n"
         "Smoked Mozz, Paneer, Feta, Farmer's Cheese, Ricotta, Mascarpone"],
        ["Certifications", "CCOF Organic, Real California Milk, rBST-free"],
        ["Key accounts", "Whole Foods, Safeway, Bi-Rite, Berkeley Bowl,\n"
         "Mollie Stone's, Draeger's, Raley's, Nugget Market"],
        ["Awards", "25-year track record: ACS 1st place (2000, 2002, 2005,\n"
         "2007, 2011), World Cheese Awards Gold (2024)"],
        ["Co-packing", "Active: Delice de la Vallee + others (comanufacturers.com)"],
        ["Owner", "Farr Hariri, ~73-75, UC Berkeley MBA '76 (Ops Research)"],
        ["Facility", "2031 2nd St, Berkeley -- Innovation Zone corridor"],
    ]
    add_mck_table(slide, data, top=Inches(1.6),
                  col_widths=[Inches(2.5), Inches(9.8)], font_size=Pt(10))

    add_gold_callout(slide,
        "37 years of continuous operation + 25 years of national awards = durable business with real craft moat",
        top=Inches(5.95))


def slide_05_succession(prs):
    slide = mck_slide(prs,
        "Farr's transition readiness is 2.5/5 -- he's in maintenance mode, not exit mode, "
        "with no formal succession plan despite being 73-75 years old",
        "Source: Agent 4 (Farr behavioral signals); Agent 2 (succession norms); family intelligence",
        slide_num=5,
        subtitle="COMPLICATION -- The succession gap")

    bullets = [
        ("FARR'S BEHAVIORAL SIGNALS", True),
        "\u2022  Still competing at ACS (2024, 2025 awards) -- invested in craft quality, not winding down",
        "\u2022  Active on Berkeley Transportation Commission -- community-engaged, not withdrawing",
        "\u2022  Co-packing still active -- developing new revenue, not harvesting",
        "\u2022  Only visible hire (2023): maintenance manager -- replacing, not expanding",
        "\u2022  Instagram 'NOW AVAILABLE IN NYC' -- Caroline may be driving distribution expansion",
        "",
        ("SUCCESSION RESEARCH (Agent 2)", True),
        "\u2022  30-40% of family businesses survive to 2nd generation (academic consensus)",
        "\u2022  Average family business succession takes 5-10 years from first conversation to completion",
        "\u2022  Most common failure mode: founder couldn't let go (cited in virtually every study)",
        "\u2022  Non-family CEO average tenure: 6.4 years vs. 17.6 for family CEOs",
        "\u2022  Non-family CEOs have 40%+ firing probability (vs. <20% for family CEOs)",
        "",
        ("IMPLICATION", True),
        "Farr may not yet recognize he needs help. The window exists but may require 1-2 years "
        "of relationship-building before Farr is ready for a serious conversation about transition.",
    ]
    add_body_bullets(slide, bullets, font_size=Pt(11))


def slide_06_nicholas(prs):
    slide = mck_slide(prs,
        "Nicholas Hariri is the critical variable: his proximity (SF), dissatisfaction (unhappy "
        "at job), and interest (showing curiosity) match the 'unexpected succession' archetype",
        "Source: Agent 10 (family dynamics); family intelligence from Nick's parents; academic research",
        slide_num=6,
        subtitle="COMPLICATION -- The Nicholas question")

    data = [
        ["Family Member", "Location", "Role", "Risk to Nick", "Key Insight"],
        ["Nicholas Hariri\n(son)", "San Francisco\n(15 min from facility)",
         "Showing interest;\nunhappy at current job", "HIGH\n(3.5/5)",
         "'Unexpected succession' archetype:\nthe child who wasn't interested, then\n"
         "suddenly is. Most dangerous for\noutside operators."],
        ["Caroline Hariri\n(daughter, older)", "New York City",
         "Marketing side gig;\ndriving NYC distribution", "LOW\n(near-term)",
         "Remote involvement limits operational\nrole. But could oppose equity dilution\n"
         "in family governance."],
        ["William Loucks\n(unknown)", "Unknown",
         "Manager per D&B\nlisting", "UNKNOWN",
         "Could be existing management layer.\nMay indicate Farr already has\n"
         "operational support."],
    ]
    add_mck_table(slide, data, top=Inches(1.6),
                  col_widths=[Inches(2.0), Inches(1.8), Inches(2.2), Inches(1.5), Inches(4.8)],
                  font_size=Pt(9))

    add_gold_callout(slide,
        "IF Nicholas commits within 2-3 years of Nick joining, Nick becomes the bridge whether "
        "he planned to be or not. Compensation must account for this risk.",
        top=Inches(5.2))

    bullets = [
        ("THREE SCENARIOS", True),
        "1. Nicholas serious + capable: Nick positions as business/ops complement (bridge role)",
        "2. Nicholas curious but uncommitted: Nick leads as primary operator (best scenario)",
        "3. Nicholas drops out: Nick becomes clear succession candidate (simplest path to equity)",
    ]
    add_body_bullets(slide, bullets, top=Inches(5.65), height=Inches(0.6), font_size=Pt(10))


def slide_07_growth(prs):
    slide = mck_slide(prs,
        "Six growth levers totaling $1.5M+ incremental revenue are available within "
        "24 months -- none require changing the product or the craft",
        "Source: Agent 3 (growth levers); Gemini research; industry benchmarks",
        slide_num=7,
        subtitle="RESOLUTION -- Growth channels")

    data = [
        ["Growth Lever", "Revenue Upside", "Capital", "Timeline", "Risk"],
        ["1. Wholesale expansion\n   (regional chains)", "$200-400K", "$20-50K",
         "6-12 mo", "Low"],
        ["2. Co-packing growth\n   (more brands)", "$200-500K", "$50-100K",
         "6-18 mo", "Low"],
        ["3. DTC / e-commerce\n   buildout", "$150-300K", "$50-100K",
         "12-18 mo", "Medium"],
        ["4. Multicultural products\n   (Paneer, Feta expansion)", "$100-200K", "$20-50K",
         "6-12 mo", "Low"],
        ["5. Burrata di Bufala\n   premium scaling", "$100-200K", "$50-100K",
         "12-24 mo", "Medium"],
        ["6. Institutional /\n   foodservice push", "$200-400K", "$20-50K",
         "6-12 mo", "Low-Med"],
        ["TOTAL", "$950K-$2.0M", "$210-450K", "12-24 mo", ""],
    ]
    add_mck_table(slide, data, top=Inches(1.6),
                  col_widths=[Inches(3.0), Inches(2.0), Inches(1.8), Inches(1.5), Inches(4.0)],
                  font_size=Pt(10))

    add_gold_callout(slide,
        "At midpoint ($1.5M incremental), Belfiore grows from $4M to $5.5M (+37%). "
        "Nick's comp effectively pays for itself through growth he enables.",
        top=Inches(5.95))


def slide_08_economics(prs, data):
    base = data["scenarios"]["base"]
    revenue = base["revenue_y1"]

    # Compute COGS, GP, EBITDA, SDE from assumptions
    cogs_pct = (base['milk_conventional_pct_rev'] + base['milk_organic_pct_rev']
                + base['milk_buffalo_pct_rev'] + base['labor_production_pct_rev']
                + base['packaging_pct_rev'] + base['utilities_ebmud_pct_rev']
                + base['distribution_pct_rev'])
    total_cogs = int(revenue * cogs_pct)
    gross_profit = revenue - total_cogs
    opex = (base['owner_comp'] + base['facility_annual'] + base['insurance_annual']
            + base['compliance_annual'] + base['marketing_annual'] + base['other_opex_annual'])
    ebitda = gross_profit - opex
    sde = ebitda + base['owner_comp']
    coo_total = int(base['coo_base_salary'] * (1 + base['coo_bonus_pct']) + base['coo_benefits_annual'])
    adj_ebitda = ebitda - coo_total
    valuation = int(sde * base['sde_multiple'])

    slide = mck_slide(prs,
        f"At ${revenue/1e6:.0f}M revenue, Belfiore generates estimated ${ebitda/1e3:.0f}K EBITDA "
        f"({ebitda/revenue:.0%} margin) with ${sde/1e3:.0f}K SDE -- the business can afford "
        f"a COO at ${coo_total/1e3:.0f}K if owner accepts a haircut",
        "Source: Agent 5 (operating economics); Agent 9 (input costs); financial model",
        slide_num=8,
        subtitle="RESOLUTION -- Business P&L (Base case)")

    left_data = [
        ["P&L Line", "% of Rev", "Est. Amount"],
        ["Revenue", "100%", f"${revenue:,}"],
        ["Milk (conventional)", f"{base['milk_conventional_pct_rev']:.0%}", f"${int(revenue*base['milk_conventional_pct_rev']):,}"],
        ["Milk (organic + buffalo)", f"{(base['milk_organic_pct_rev']+base['milk_buffalo_pct_rev']):.0%}",
         f"${int(revenue*(base['milk_organic_pct_rev']+base['milk_buffalo_pct_rev'])):,}"],
        ["Production labor", f"{base['labor_production_pct_rev']:.0%}", f"${int(revenue*base['labor_production_pct_rev']):,}"],
        ["Packaging + utilities", f"{(base['packaging_pct_rev']+base['utilities_ebmud_pct_rev']):.1%}",
         f"${int(revenue*(base['packaging_pct_rev']+base['utilities_ebmud_pct_rev'])):,}"],
        ["Distribution", f"{base['distribution_pct_rev']:.1%}", f"${int(revenue*base['distribution_pct_rev']):,}"],
        ["Total COGS", f"{cogs_pct:.1%}", f"${total_cogs:,}"],
        ["Gross Profit", f"{(1-cogs_pct):.1%}", f"${gross_profit:,}"],
        ["Operating expenses", "", f"${opex:,}"],
        ["EBITDA", f"~{ebitda/revenue:.0%}", f"~${ebitda:,}"],
        ["SDE (incl. owner comp)", "", f"~${sde:,}"],
    ]
    add_mck_table(slide, left_data, top=Inches(1.6), left=CONTENT_L, width=Inches(6.5),
                  col_widths=[Inches(3.0), Inches(1.5), Inches(2.0)], font_size=Pt(10))

    right_data = [
        ["With COO (Nick)", "Impact"],
        ["COO total comp", f"~${coo_total:,}"],
        ["Adjusted EBITDA", f"~${adj_ebitda:,}"],
        ["Owner comp reduction", f"~${coo_total:,}/yr"],
        ["Business valuation", f"~${valuation:,}"],
        [f"(SDE x {base['sde_multiple']:.1f})", "(at current SDE)"],
    ]
    add_mck_table(slide, right_data, top=Inches(1.6), left=Inches(7.3), width=Inches(5.5),
                  col_widths=[Inches(3.0), Inches(2.5)], font_size=Pt(10))

    add_gold_callout(slide,
        f"Farr takes a ~${coo_total:,} annual haircut to his income. This only works if he values "
        "freed time + growth + succession clarity more than that income.",
        top=Inches(5.95))


def slide_09_personal_decision(prs, data):
    base = data["scenarios"]["base"]
    revenue = base["revenue_y1"]

    # Compute SDE from assumptions (same logic as slide_08)
    cogs_pct = (base['milk_conventional_pct_rev'] + base['milk_organic_pct_rev']
                + base['milk_buffalo_pct_rev'] + base['labor_production_pct_rev']
                + base['packaging_pct_rev'] + base['utilities_ebmud_pct_rev']
                + base['distribution_pct_rev'])
    gross_profit = revenue - int(revenue * cogs_pct)
    opex = (base['owner_comp'] + base['facility_annual'] + base['insurance_annual']
            + base['compliance_annual'] + base['marketing_annual'] + base['other_opex_annual'])
    ebitda = gross_profit - opex
    sde = ebitda + base['owner_comp']

    coo_total = int(base["coo_base_salary"] * (1 + base["coo_bonus_pct"]) + base["coo_benefits_annual"])
    trad_y1 = base["traditional_career_comp_y1"]
    trad_g = base["traditional_career_growth"]

    slide = mck_slide(prs,
        "The personal decision hinges on equity: without it, Nick sacrifices $200K+ "
        "over 5 years vs. traditional career; with 5% annual vesting, he breaks even by Year 4",
        "Source: Financial model (Personal Decision sheet); comp benchmarks (Agent 1)",
        slide_num=9,
        subtitle="RESOLUTION -- Nick's 5-year path comparison")

    # Build 5-year comparison
    trad_cum = 0
    bridge_cum = 0
    equity_cum = 0
    mult = base["sde_multiple"]

    path_data = [["Year", "Path A:\nTraditional", "Path B:\nBridge (no equity)", "Path C:\nWith 5% vesting",
                  "Equity Value\n(Path C)"]]

    for yr in range(1, 6):
        trad_yr = int(trad_y1 * (1 + trad_g) ** (yr - 1))
        trad_cum += trad_yr
        bridge_cum += coo_total
        equity_cum += coo_total
        eq_pct = base["equity_vest_pct_per_year"] * yr
        eq_val = int(eq_pct * sde * mult)

        path_data.append([
            f"Year {yr}",
            f"${trad_cum:,}",
            f"${bridge_cum:,}",
            f"${equity_cum + eq_val:,}",
            f"${eq_val:,}\n({eq_pct:.0%} vested)",
        ])

    path_data.append([
        "Opportunity\nCost vs. A",
        "--",
        f"-${trad_cum - bridge_cum:,}",
        f"{'+'if (equity_cum + eq_val - trad_cum) >= 0 else '-'}${abs(equity_cum + eq_val - trad_cum):,}",
        "",
    ])

    add_mck_table(slide, path_data, top=Inches(1.6),
                  col_widths=[Inches(1.5), Inches(2.5), Inches(2.8), Inches(2.8), Inches(2.7)],
                  font_size=Pt(9))

    add_gold_callout(slide,
        "KEY QUESTION: At what equity vesting rate does Path C break even with Path A? "
        "Change assumptions in the Excel model to find the answer.",
        top=Inches(5.95))


def slide_10_entry_strategy(prs):
    slide = mck_slide(prs,
        "The entry strategy is a 20-week phased relationship arc from social introduction "
        "to formal diligence -- with 6 kill criteria committed to in advance",
        "Source: Diligence plan; Agent 7 (case studies); HBR family business research",
        slide_num=10,
        subtitle="RESOLUTION -- How Nick gets in")

    data = [
        ["Phase", "Weeks", "Objective", "Kill Signal"],
        ["1. Relationship Entry", "1-2",
         "Parents facilitate introduction.\nFirst coffee/dinner -- listen, be curious.",
         "Farr uninterested in deeper\nconnection."],
        ["2. Relationship Deepening", "3-6",
         "Facility visit (if invited). Business\nconversation only if Farr opens door.",
         "Kill 1: No interest in help.\nKill 5: Nicholas is the heir."],
        ["3. Structured Exploration", "6-10",
         "Compensation reality check. Role\ndefinition. Family alignment.",
         "Kill 2: No equity path.\nKill 3: Comp floor unmet."],
        ["4. Due Diligence", "10-16",
         "Financial review. Operational\nassessment. Legal structuring.",
         "Kill 4: Business declining.\nKill 6: Skills don't fit."],
        ["5. Decision", "16-20",
         "Internal decision memo. Go/no-go.\nCommitment or graceful exit.",
         "All accumulated evidence\nweighed against kill criteria."],
    ]
    add_mck_table(slide, data, top=Inches(1.6),
                  col_widths=[Inches(2.5), Inches(1.0), Inches(4.5), Inches(4.3)],
                  font_size=Pt(10))

    add_gold_callout(slide,
        "GOVERNING PRINCIPLE: 'Would Farr feel respected, not threatened, if he saw this?' "
        "-- every action, artifact, and conversation must pass this test.",
        top=Inches(5.95))


def slide_11_risks(prs):
    slide = mck_slide(prs,
        "Eight risks identified -- the two highest-impact risks (no equity path, "
        "Nicholas preemption) are mitigated by pre-commitment decisions, not hope",
        "Source: Risk matrix from pre-call briefing; case studies; academic research",
        slide_num=11,
        subtitle="RESOLUTION -- Risk register")

    data = [
        ["Risk", "Likelihood", "Impact", "Mitigation"],
        ["Farr doesn't want help", "40%", "FATAL",
         "Kill Criteria #1. Read signals Meeting 1-2."],
        ["Equity never materializes", "50%", "FATAL",
         "Written agreement within 6 months or exit."],
        ["Nicholas enters full-time\nwithin 2 years", "35%", "HIGH",
         "Negotiate severance + phantom equity\nthat vests on departure."],
        ["Below-market comp\nwith no upside", "30%", "HIGH",
         "Define floor ($120K+) before any\nconversation. Walk if not met."],
        ["Family opposition", "35%", "HIGH",
         "Farr must socialize with family.\nDon't go around him."],
        ["Cultural mismatch\n(McKinsey vs. artisan)", "30%", "MED",
         "Lead with curiosity, not frameworks.\nFirst 6 months = learning mode."],
        ["Business in harvest mode", "20%", "HIGH",
         "Facility visit + account check in Phase 2."],
        ["Berkeley property value\nexceeds business value", "25%", "MED",
         "Understand lease vs. own.\nIf owned, land is the real asset."],
    ]
    add_mck_table(slide, data, top=Inches(1.6),
                  col_widths=[Inches(2.8), Inches(1.3), Inches(1.2), Inches(7.0)],
                  font_size=Pt(10))


def slide_12_next_steps(prs):
    slide = mck_slide(prs,
        "Three immediate actions, three pre-commitment decisions, and one governing "
        "truth -- proceed with the introduction regardless of the operating entry outcome",
        "Source: Internal analysis",
        slide_num=12)

    bullets = [
        ("IMMEDIATE ACTIONS", True),
        "1. Nick's parents facilitate a social introduction to Farr. Framing: curiosity about the craft, not a job pitch.",
        "2. Nick defines his compensation floor, equity timeline, and walk-away conditions -- privately, before any conversation.",
        "3. Nick answers: 'Am I willing to be a bridge CEO with no equity?' If yes, comp must be market-rate ($150K+).",
        "",
        ("PRE-COMMITMENT DECISIONS", True),
        "\u2022  Compensation floor: $___K total (suggest $120K minimum including benefits)",
        "\u2022  Equity timeline: Written agreement within ___ months of joining (suggest 6) or exit",
        "\u2022  Walk-away scenario: If I haven't achieved ___ within ___ years, I leave. Define now.",
        "",
        ("THE GOVERNING TRUTH", True),
        "The relationship with Farr is worth building regardless of the operating entry outcome. "
        "Farr is a 35-year Bay Area food entrepreneur with deep industry knowledge, a Berkeley MBA, "
        "and a network of suppliers, retailers, and fellow artisans. Even if the operating entry "
        "doesn't work, the mentorship and network are valuable.",
        "",
        "But go in with eyes open: equity is earned at the negotiating table, not at the cheese vat. "
        "And zero outside operators at this scale have earned it through sweat equity alone.",
    ]
    add_body_bullets(slide, bullets, font_size=Pt(11))

    add_gold_callout(slide,
        "The most honest framing: a 20% shot at life-changing ownership, a 40% shot at "
        "valuable experience, and a 25% shot at a good mentor. Proceed accordingly.",
        top=Inches(5.95))


# ═════════════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════════════
def main():
    os.makedirs(OUT_DIR, exist_ok=True)

    with open(ASSUMPTIONS_FILE, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_cover(prs)
    slide_02_exec_summary(prs)
    slide_03_market(prs)
    slide_04_company(prs)
    slide_05_succession(prs)
    slide_06_nicholas(prs)
    slide_07_growth(prs)
    slide_08_economics(prs, data)
    slide_09_personal_decision(prs, data)
    slide_10_entry_strategy(prs)
    slide_11_risks(prs)
    slide_12_next_steps(prs)

    prs.save(OUT_FILE)
    print(f"OK: {OUT_FILE} ({os.path.getsize(OUT_FILE):,} bytes)")
    print(f"12 slides: Cover, Exec Summary, Market, Company, Succession, Nicholas,")
    print(f"           Growth, Economics, Personal Decision, Entry Strategy, Risks, Next Steps")


if __name__ == "__main__":
    main()
