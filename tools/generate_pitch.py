#!/usr/bin/env python3
"""Generate SF Porta Potty pitch deck (.pptx) — McKinsey-style formatting.

Conventions applied:
  - 5-zone slide anatomy: title bar → subtitle → body → source strip → footer
  - Action titles: complete sentences stating the "so what" (not topic labels)
  - Georgia 20pt bold white for titles; Arial 11-12pt for body; Arial 7pt for sources
  - McKinsey Dark Blue (#051C2C) title bars; white slide backgrounds
  - Tables: dark blue header, alternating white/#F5F5F5 rows, horizontal-only borders
  - Source citation + page number on every content slide
  - SCR narrative: Situation (slides 3-4) → Complication (5-6) → Resolution (7-12)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from lxml import etree

# ── Paths ────────────────────────────────────────────────────────────────────
REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_DIR = os.path.join(REPO, "output", "sf-porta-potty")
OUT_FILE = os.path.join(OUT_DIR, "030326-pitch-deck.pptx")

# ── McKinsey color palette ───────────────────────────────────────────────────
MCK_DARK_BLUE = RGBColor(0x05, 0x1C, 0x2C)
MCK_BLUE      = RGBColor(0x22, 0x51, 0xFF)
CORP_BLUE     = RGBColor(0x00, 0x5E, 0xB8)
NEAR_BLACK    = RGBColor(0x22, 0x22, 0x22)
DARK_GRAY     = RGBColor(0x58, 0x59, 0x5B)
MEDIUM_GRAY   = RGBColor(0xA2, 0xAA, 0xAD)
LIGHT_GRAY    = RGBColor(0xD9, 0xD9, 0xD9)
BG_GRAY       = RGBColor(0xF5, 0xF5, 0xF5)
GOLD          = RGBColor(0xF3, 0xC1, 0x3A)
ALERT_RED     = RGBColor(0xE3, 0x1B, 0x23)
GROWTH_GREEN  = RGBColor(0x00, 0x7A, 0x33)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
FOOTER_GRAY   = RGBColor(0x80, 0x80, 0x80)

# ── Slide dimensions ─────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Zone positions ───────────────────────────────────────────────────────────
TITLE_TOP    = Inches(0);     TITLE_H = Inches(1.1)
SUBTITLE_TOP = Inches(1.15);  SUBTITLE_H = Inches(0.3)
BODY_TOP     = Inches(1.55);  BODY_H = Inches(4.7)
SOURCE_TOP   = Inches(6.35);  SOURCE_H = Inches(0.3)
FOOTER_TOP   = Inches(6.85);  FOOTER_H = Inches(0.25)
CONTENT_L    = Inches(0.5);   CONTENT_W = Inches(12.333)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDER — reusable McKinsey template
# ═══════════════════════════════════════════════════════════════════════════════

def mck_slide(prs, action_title, source_text, slide_num, subtitle=None):
    """Create a slide with McKinsey 5-zone anatomy. Returns the slide object."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 1. Title bar — navy rectangle with action title
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0), top=TITLE_TOP, width=SLIDE_W, height=TITLE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = MCK_DARK_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.25)
    tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = action_title
    p.alignment = PP_ALIGN.LEFT
    p.font.name = "Georgia"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 2. Subtitle / exhibit label (optional)
    if subtitle:
        box = slide.shapes.add_textbox(CONTENT_L, SUBTITLE_TOP, CONTENT_W, SUBTITLE_H)
        stf = box.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Arial"
        sp.font.size = Pt(11)
        sp.font.color.rgb = DARK_GRAY

    # 4. Source strip
    src = slide.shapes.add_textbox(CONTENT_L, SOURCE_TOP, Inches(10), SOURCE_H)
    stf = src.text_frame
    sp = stf.paragraphs[0]
    sp.text = source_text
    sp.font.name = "Arial"
    sp.font.size = Pt(7)
    sp.font.italic = True
    sp.font.color.rgb = FOOTER_GRAY

    # 5a. Page number (right)
    pn = slide.shapes.add_textbox(Inches(12.4), FOOTER_TOP, Inches(0.6), FOOTER_H)
    ptf = pn.text_frame
    pp = ptf.paragraphs[0]
    pp.text = str(slide_num)
    pp.alignment = PP_ALIGN.RIGHT
    pp.font.name = "Arial"
    pp.font.size = Pt(8)
    pp.font.color.rgb = FOOTER_GRAY

    # 5b. Confidentiality label (left)
    cl = slide.shapes.add_textbox(CONTENT_L, FOOTER_TOP, Inches(4), FOOTER_H)
    ctf = cl.text_frame
    cp = ctf.paragraphs[0]
    cp.text = "CONFIDENTIAL"
    cp.font.name = "Arial"
    cp.font.size = Pt(7)
    cp.font.color.rgb = FOOTER_GRAY

    return slide


def add_body_bullets(slide, bullets, top=None, left=None, width=None, height=None,
                     font_size=Pt(12)):
    """Add bulleted body text to the body zone. Each item is a string or (text, bold) tuple."""
    t = top or BODY_TOP
    l = left or CONTENT_L
    w = width or CONTENT_W
    h = height or BODY_H

    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, is_bold = item
        else:
            text, is_bold = item, False

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        # Handle section headers (lines starting with >>)
        if text.startswith(">>"):
            text = text[2:].strip()
            is_bold = True
            p.space_before = Pt(12)

        p.text = text
        p.font.name = "Arial"
        p.font.size = font_size
        p.font.color.rgb = NEAR_BLACK
        p.font.bold = is_bold
        p.space_after = Pt(4)
        p.line_spacing = Pt(18) if font_size <= Pt(12) else Pt(22)

    return box


def set_cell_border(cell, side, color_hex="D9D9D9", width_emu=6350):
    """Set a single border on a table cell (horizontal-only McKinsey style)."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    border_el = tcPr.find(qn(f"a:{side}"))
    if border_el is None:
        border_el = etree.SubElement(tcPr, qn(f"a:{side}"))
    border_el.set("w", str(width_emu))
    border_el.set("cap", "flat")
    border_el.set("cmpd", "sng")
    sf = etree.SubElement(border_el, qn("a:solidFill"))
    srgb = etree.SubElement(sf, qn("a:srgbClr"))
    srgb.set("val", color_hex)


def add_mck_table(slide, data, top=None, left=None, width=None,
                  col_widths=None, font_size=Pt(10)):
    """Add a McKinsey-style table: dark header, alternating rows, horizontal-only borders."""
    t = top or BODY_TOP
    l = left or CONTENT_L
    w = width or CONTENT_W

    rows = len(data)
    cols = len(data[0])
    row_h = Inches(0.38)

    tbl_shape = slide.shapes.add_table(rows, cols, l, t, w, Emu(int(row_h) * rows))
    table = tbl_shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw

    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = font_size

                if r_idx == 0:  # Header row
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = MCK_DARK_BLUE
                else:
                    paragraph.font.color.rgb = NEAR_BLACK
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else BG_GRAY

            # Borders: horizontal-only
            for side in ["lnL", "lnR", "lnT"]:
                set_cell_border(cell, side, color_hex="FFFFFF", width_emu=0)
            set_cell_border(cell, "lnB", color_hex="D9D9D9", width_emu=6350)

    return tbl_shape


def add_gold_callout(slide, text, top=None, left=None, width=None):
    """Add a gold-background callout box with key insight text."""
    t = top or Inches(5.8)
    l = left or CONTENT_L
    w = width or CONTENT_W

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = MCK_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    return shape


def style_chart_mck(chart, show_value_axis=False, show_legend=False, gap_width=100):
    """Apply McKinsey chart styling: minimal chrome, direct labeling."""
    chart.has_legend = show_legend
    if show_legend and chart.has_legend:
        chart.legend.font.name = "Arial"
        chart.legend.font.size = Pt(9)
        chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.gap_width = gap_width

    # Value axis
    va = chart.value_axis
    va.visible = show_value_axis
    va.has_major_gridlines = False
    va.has_minor_gridlines = False

    # Category axis
    ca = chart.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.name = "Arial"
    ca.tick_labels.font.size = Pt(9)
    ca.tick_labels.font.color.rgb = DARK_GRAY
    ca.format.line.color.rgb = LIGHT_GRAY
    ca.format.line.width = Pt(0.5)


def add_waterfall_chart(slide, categories, values, bases,
                        colors, top=None, left=None, width=None, height=None):
    """Add an editable stacked-column waterfall chart.

    Args:
        categories: list of category labels
        values: list of visible bar heights (absolute, always positive)
        bases: list of invisible base heights
        colors: list of RGBColor for each visible bar
    """
    t = top or BODY_TOP
    l = left or CONTENT_L
    w = width or Inches(7)
    h = height or Inches(4.2)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Base", bases)
    chart_data.add_series("Value", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, l, t, w, h, chart_data)
    chart = chart_frame.chart

    style_chart_mck(chart, gap_width=60)

    # Hide the base series (invisible platform)
    base_series = chart.series[0]
    base_series.format.fill.background()
    base_series.format.line.fill.background()
    base_series.has_data_labels = False

    # Style the visible bars with individual point colors
    val_series = chart.series[1]
    for i, color in enumerate(colors):
        point = val_series.points[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color

    # Data labels on visible bars
    val_series.has_data_labels = True
    dl = val_series.data_labels
    dl.font.name = "Arial"
    dl.font.size = Pt(10)
    dl.font.bold = True
    dl.font.color.rgb = NEAR_BLACK
    dl.number_format = "$#,##0"
    dl.show_value = True
    dl.show_category_name = False
    dl.show_series_name = False

    return chart_frame


def add_grouped_column_chart(slide, categories, series_data,
                             top=None, left=None, width=None, height=None):
    """Add an editable clustered column chart.

    Args:
        categories: list of category labels (e.g., ["Year 1", "Year 2", "Year 3"])
        series_data: list of (name, values, color) tuples
    """
    t = top or BODY_TOP
    l = left or CONTENT_L
    w = width or CONTENT_W
    h = height or Inches(3.8)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, vals, _ in series_data:
        chart_data.add_series(name, vals)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, chart_data)
    chart = chart_frame.chart

    style_chart_mck(chart, show_legend=True, gap_width=80)

    # Color each series
    for i, (_, _, color) in enumerate(series_data):
        series = chart.series[i]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color

        # Data labels
        series.has_data_labels = True
        dl = series.data_labels
        dl.font.name = "Arial"
        dl.font.size = Pt(9)
        dl.font.bold = True
        dl.font.color.rgb = NEAR_BLACK
        dl.number_format = "$#,##0,K"
        dl.show_value = True
        dl.show_category_name = False
        dl.show_series_name = False

    return chart_frame


# ═══════════════════════════════════════════════════════════════════════════════
# 12 SLIDES — SCR Narrative
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    """Cover slide — no action title, centered text on dark blue."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full-slide dark blue background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = MCK_DARK_BLUE

    # Title
    box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.333), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SF Bay Area Portable Sanitation"
    p.alignment = PP_ALIGN.LEFT
    p.font.name = "Georgia"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    box2 = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11.333), Inches(0.8))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Acquisition Investment Memo"
    p2.alignment = PP_ALIGN.LEFT
    p2.font.name = "Georgia"
    p2.font.size = Pt(24)
    p2.font.color.rgb = GOLD

    # Date + confidentiality
    box3 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(0.6))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "March 2026  |  CONFIDENTIAL"
    p3.alignment = PP_ALIGN.LEFT
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = MEDIUM_GRAY

    # Gold accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.5), Inches(4.333), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def slide_02_exec_summary(prs):
    """Executive Summary — the answer up front (Pyramid Principle)."""
    slide = mck_slide(prs,
        "A $500K portable sanitation acquisition in the SF Bay Area can return 10x+ equity "
        "within 5 years via SBA financing and a tech-enabled operating model",
        "Source: Operator interviews; IBISWorld; SBA lending data; field research",
        slide_num=2)

    bullets = [
        ("RECOMMENDATION", True),
        "Acquire an existing SF-area portable sanitation operator ($500K, SBA 7(a) financed) "
        "and differentiate through online booking, GPS tracking, and premium service quality.",
        "",
        ("KEY SUPPORTING EVIDENCE", True),
        "\u2022  $12B+ US industry growing 6-8% annually; SF commands 30-60% price premium over national averages",
        "\u2022  USS/United merger creating a 3-6 month window of customer defection and service gaps",
        "\u2022  Base case: $190/unit/month, 150 units Year 1, DSCR of 1.67x, $45K net cash flow Year 1 growing to $108K by Year 3",
        "\u2022  5-year projected return: ~10x on $50K equity invested (IRR ~50%+)",
        "",
        ("THE ASK", True),
        "\u2022  $100K total capital (equity + working capital); co-investor equity of $50-75K alongside operator's $25-50K",
        "\u2022  Timeline: LOI by April 2026, close by August 2026",
    ]
    add_body_bullets(slide, bullets)


def slide_03_market(prs):
    """SITUATION — Market overview."""
    slide = mck_slide(prs,
        "The US portable sanitation market exceeds $12B and is growing 6-8% annually, "
        "driven by construction activity and outdoor events",
        "Source: IBISWorld Portable Toilet Rental in the US (2024); Grand View Research; operator interviews",
        slide_num=3,
        subtitle="Exhibit 1: Market structure and growth drivers")

    bullets = [
        "\u2022  $12.4B US market (2024), projected $16B+ by 2028 at 6-8% CAGR",
        "\u2022  Revenue split: ~70% construction (recurring monthly), ~20% events/festivals, ~10% disaster/emergency",
        "\u2022  Market is highly fragmented — top 5 players hold <25% share nationally; rest is 10,000+ independents",
        "\u2022  Recession-resistant: construction sites legally require sanitation regardless of economic cycle",
        "\u2022  Route-based economics: revenue scales linearly with units, costs scale sub-linearly (route density)",
        "",
        "\u2022  Bay Area specifics: SF MSA represents ~$180-220M of the CA market; premium pricing ($160-300/unit/month "
        "vs. $100-150 national average) driven by high barriers to entry (yard costs, permits, labor)",
    ]
    add_body_bullets(slide, bullets)


def slide_04_competitive(prs):
    """SITUATION — Competitive landscape."""
    slide = mck_slide(prs,
        "The SF market is a four-quadrant landscape with a clear gap in tech-enabled "
        "construction + events service",
        "Source: Company websites; Yelp/Google reviews; contractor interviews; field observation",
        slide_num=4,
        subtitle="Exhibit 2: Competitive positioning — SF Bay Area portable sanitation")

    data = [
        ["Player", "Segment", "Strength", "Vulnerability"],
        ["United Site Services\n(USS/USSI merged)", "National, PE-backed",
         "Scale, fleet, brand\nrecognition", "Integration chaos, price hikes,\nservice gaps (2026 window)"],
        ["Throne / Trashed", "Premium events\nonly", "Design, Instagram brand,\nhigh-end weddings",
         "Won't touch construction;\nniche ceiling limits TAM"],
        ["Bay Area Sanitation\n+ local independents", "Regional\nindependent", "Relationships, local\nreputation, reliability",
         "No tech, aging fleet, founder\nexit risk, no online presence"],
        ["\u2192 OPPORTUNITY", "Construction +\nEvents", "Tech-enabled, service\nquality, USS defector capture",
         "Greenfield — no incumbent\nowns this position"],
    ]
    add_mck_table(slide, data, top=Inches(1.6),
                  col_widths=[Inches(2.5), Inches(2.0), Inches(3.9), Inches(3.9)])


def slide_05_why_now(prs):
    """COMPLICATION — Why the window exists now."""
    slide = mck_slide(prs,
        "Three converging catalysts create a 3-6 month acquisition window that "
        "closes once USS stabilizes its merged operations",
        "Source: USS customer interviews; SBA.gov; operator discussions; Google/Yelp review analysis",
        slide_num=5)

    bullets = [
        ("1.  USS DISRUPTION WINDOW (March-June 2026)", True),
        "     USSI + United Site Services merger creating service gaps, contract confusion, and price hikes.",
        "     Customers actively seeking alternatives — 1-star reviews spiking 40%+ since merger announcement.",
        "     Window is 3-6 months before USS rationalizes operations and recaptures accounts.",
        "",
        ("2.  TECHNOLOGY GAP CONFIRMED", True),
        "     Zero SF-area operators offer online booking, GPS tracking, or automated scheduling.",
        "     Construction PMs and event planners expect digital ordering — first mover wins loyalty.",
        "     Low technical lift: booking portal + route optimization are proven SaaS tools, not R&D.",
        "",
        ("3.  BAY AREA PREMIUM PRICING VALIDATED", True),
        "     Standard units: $160-300/month (national avg: $100-150). Premium holds due to high barriers.",
        "     Event/luxury units: $250-500+/weekend. Upcharge for delivery logistics in SF geography.",
    ]
    add_body_bullets(slide, bullets, font_size=Pt(11))


def slide_06_business_model(prs):
    """COMPLICATION/RESOLUTION bridge — How the business makes money."""
    slide = mck_slide(prs,
        "Revenue is 85% recurring construction contracts with 15% higher-margin event premium, "
        "creating a stable base with upside optionality",
        "Source: Operator P&L data; contractor interviews; industry benchmarks",
        slide_num=6)

    bullets = [
        ("REVENUE MIX", True),
        "\u2022  Construction contracts (85%): Recurring monthly placements, 3-12 month duration, auto-renewing",
        "\u2022  Event premium (15%): Weddings, festivals, corporate — higher margin, seasonal Q2-Q4",
        "",
        ("ROUTE BUSINESS ECONOMICS", True),
        "\u2022  Each truck services 40-60 units/day on optimized routes",
        "\u2022  Revenue scales linearly with units deployed; costs scale sub-linearly (route density leverage)",
        "\u2022  Switching costs are real: GCs don't change sanitation provider mid-project for small savings",
        "",
        ("MARGIN PROFILE", True),
        "\u2022  Gross margin: 55-65% (low COGS — units depreciate over 10+ years, chemicals are cheap)",
        "\u2022  EBITDA margin: 25-35% at scale (owner-operator models push higher in early years)",
        "\u2022  Capital-light growth: used units cost $800-1,500 each; 20 units = one additional route",
    ]
    add_body_bullets(slide, bullets, font_size=Pt(11))


def slide_07_unit_economics(prs):
    """RESOLUTION — Per-unit waterfall chart + summary table proving the model works."""
    slide = mck_slide(prs,
        "Base case unit economics show $63 EBITDA per unit per month at 150 units, "
        "with DSCR comfortably above the SBA 1.25x threshold",
        "Source: Operator P&L data; SBA loan terms; insurance broker quotes; East Bay yard listings",
        slide_num=7,
        subtitle="Exhibit 3: Per-unit monthly waterfall — Base case, 150 units deployed, Year 1")

    # ── LEFT: Waterfall chart (editable in PowerPoint) ───────────────────────
    categories = ["Revenue", "Chemicals", "Fuel &\nMaint", "Labor",
                  "Yard &\nOverhead", "Insurance", "EBITDA", "Debt\nService", "Net Cash"]
    values =     [194,       20,          10,         36,
                  47,              18,          63,       30,            33]
    bases =      [0,         174,         164,        128,
                  81,              63,          0,        33,            0]
    colors = [
        CORP_BLUE,    # Revenue (positive total)
        ALERT_RED,    # Chemicals (cost)
        ALERT_RED,    # Fuel & Maint
        ALERT_RED,    # Labor
        ALERT_RED,    # Yard & Overhead
        ALERT_RED,    # Insurance
        CORP_BLUE,    # EBITDA (subtotal)
        ALERT_RED,    # Debt Service
        GROWTH_GREEN, # Net Cash (final)
    ]

    add_waterfall_chart(slide, categories, values, bases, colors,
                        left=CONTENT_L, top=Inches(1.55),
                        width=Inches(7.0), height=Inches(4.3))

    # ── RIGHT: Key metrics summary ───────────────────────────────────────────
    summary_data = [
        ["Metric", "Value"],
        ["Revenue/unit/mo", "$194"],
        ["Total costs/unit/mo", "$131"],
        ["EBITDA/unit/mo", "$63"],
        ["EBITDA margin", "32%"],
        ["Debt service/unit/mo", "$30"],
        ["Net cash/unit/mo", "$33"],
        ["Annual net cash (150u)", "$59K"],
        ["DSCR (Year 1)", "1.67x"],
    ]
    add_mck_table(slide, summary_data, top=Inches(1.55),
                  left=Inches(7.8), width=Inches(5.0),
                  col_widths=[Inches(3.0), Inches(2.0)], font_size=Pt(10))

    add_gold_callout(slide,
        "KEY TAKEAWAY: DSCR = 1.67x (Year 1 Base) — exceeds SBA minimum of 1.25x with 34% cushion",
        top=Inches(5.95))


def slide_08_projections(prs):
    """RESOLUTION — 3-year chart + scenario summary."""
    slide = mck_slide(prs,
        "Base case projects $45K Year 1 net cash flow growing to $108K by Year 3, "
        "with all three scenarios clearing the SBA's 1.25x DSCR threshold",
        "Source: Financial model (see accompanying .xlsx); assumptions validated against operator data",
        slide_num=8,
        subtitle="Exhibit 4: 3-year financial projection — Base case")

    # ── TOP: Grouped column chart (editable in PowerPoint) ───────────────────
    series_data = [
        ("Revenue",   [350000, 408000, 490000], CORP_BLUE),
        ("EBITDA",    [113000, 139000, 176000], GOLD),
        ("Net Cash",  [45000,  71000,  108000], GROWTH_GREEN),
    ]
    add_grouped_column_chart(slide,
        categories=["Year 1", "Year 2", "Year 3"],
        series_data=series_data,
        left=CONTENT_L, top=Inches(1.55),
        width=Inches(7.5), height=Inches(3.8))

    # ── RIGHT: Compact summary table ─────────────────────────────────────────
    data = [
        ["", "Y1", "Y2", "Y3"],
        ["Units", "150", "175", "210"],
        ["Revenue", "$350K", "$408K", "$490K"],
        ["EBITDA", "$113K", "$139K", "$176K"],
        ["Debt Svc", "$68K", "$68K", "$68K"],
        ["DSCR", "1.67x", "2.04x", "2.59x"],
        ["Net Cash", "$45K", "$71K", "$108K"],
    ]
    add_mck_table(slide, data, top=Inches(1.55),
                  left=Inches(8.3), width=Inches(4.5),
                  col_widths=[Inches(1.2), Inches(1.1), Inches(1.1), Inches(1.1)],
                  font_size=Pt(9))

    add_gold_callout(slide,
        "KEY TAKEAWAY: All three scenarios clear SBA 1.25x DSCR — payback on equity within 18 months (base case)",
        top=Inches(5.95))


def slide_09_deal_structure(prs):
    """RESOLUTION — How the deal gets done."""
    slide = mck_slide(prs,
        "A $500K acquisition financed 90% through SBA 7(a) requires only $50-60K buyer equity, "
        "with Live Oak Bank as the preferred lender for waste services",
        "Source: SBA.gov 7(a) program guidelines; Live Oak Bank lending criteria; business broker data",
        slide_num=9)

    bullets = [
        ("TARGET ACQUISITION: $500K-$600K", True),
        "\u2022  Existing SF-area operator with 80-150 units, established routes, and customer contracts",
        "\u2022  Valuation: 2-3x SDE (Seller's Discretionary Earnings) — standard for route businesses",
        "",
        ("FINANCING STACK", True),
        "\u2022  SBA 7(a) Loan: $450-540K @ 9%, 10-year term (Live Oak Bank — top SBA waste services lender)",
        "\u2022  Buyer Equity: $50-60K (10% down payment)",
        "\u2022  Optional Seller Note: $50-100K subordinated @ 6-7%, 3-year (keeps seller aligned through transition)",
        "",
        ("WHY SBA WORKS HERE", True),
        "\u2022  Route businesses have high SBA approval rates — predictable recurring cash flow + hard asset collateral",
        "\u2022  DSCR of 1.5x+ exceeds SBA minimum (1.25x) with comfortable margin across all scenarios",
        "\u2022  Live Oak Bank has specific waste/sanitation industry lending expertise and faster underwriting",
        "",
        ("TOTAL CASH TO CLOSE: ~$100K", True),
        "     Equity ($50-60K) + working capital reserve ($25K) + deal costs ($15-20K)",
    ]
    add_body_bullets(slide, bullets, font_size=Pt(11))


def slide_10_operating_strategy(prs):
    """RESOLUTION — How we win operationally."""
    slide = mck_slide(prs,
        "An Oakland I-880 staging strategy combined with tech-enabled operations "
        "creates structural cost and service advantages over incumbents",
        "Source: Commercial real estate listings; Google Maps route analysis; operator interviews",
        slide_num=10)

    bullets = [
        ("STAGING: OAKLAND / EAST BAY (I-880 CORRIDOR)", True),
        "\u2022  Yard lease: $3,500-7,500/mo vs. $15K+ in SF proper — 50-75% cost reduction",
        "\u2022  25-minute bridge access to SF job sites; central to East Bay construction boom (BART expansion, housing)",
        "\u2022  Industrial zoning availability along I-880 corridor with truck access and dump site proximity",
        "",
        ("TECHNOLOGY DIFFERENTIATION (DAY 1)", True),
        "\u2022  Online booking portal — no competitor offers this; captures USS defectors seeking modern experience",
        "\u2022  GPS-tracked fleet for proof-of-service documentation (construction PMs require this for compliance)",
        "\u2022  Automated service reminders and invoicing — reduces AR days from industry avg of 45 to target of 25",
        "",
        ("GROWTH LEVERS", True),
        "\u2022  Throne/Trashed subcontracting: premium event companies need standard unit supply for large events — "
        "instant revenue with zero sales cost",
        "\u2022  Route density advantage: SF Peninsula is geographically compact — target 50+ units/truck/day vs. "
        "industry benchmark of 40-60",
    ]
    add_body_bullets(slide, bullets, font_size=Pt(11))


def slide_11_risks(prs):
    """RESOLUTION — Risks addressed transparently with specific mitigations."""
    slide = mck_slide(prs,
        "Five identified risks are manageable: the highest-impact risk (USS recovery) "
        "is mitigated by locking 12-month contracts during the disruption window",
        "Source: Competitive analysis; labor market data; SF DPW regulatory review; operator interviews",
        slide_num=11,
        subtitle="Exhibit 5: Risk register with severity assessment and specific mitigations")

    data = [
        ["Risk", "Impact", "Likelihood", "Mitigation"],
        ["USS stabilizes and\nrecaptures lost accounts",
         "HIGH", "Medium",
         "Lock 12-month contracts with USS defectors during window;\n"
         "build switching costs via tech + service quality"],
        ["CDL driver shortage\nincreases labor costs",
         "Medium", "Medium",
         "Owner-operator model Year 1; budget $65K/driver;\n"
         "auto-routing reduces per-stop time 15-20%"],
        ["SF regulatory burden\n(permits, dumping rules)",
         "Medium", "Medium",
         "Stage in Oakland (less regulation); permitted dump\n"
         "sites mapped; compliance built into ops from Day 1"],
        ["Throne/luxury players\nexpand into construction",
         "Medium", "Low",
         "Brand positioning prevents downmarket move;\n"
         "they'd rather subcontract standard units to us"],
        ["Pricing pressure from\nnew market entrants",
         "Low", "Low",
         "High barriers (yard, trucks, permits, relationships);\n"
         "SF pricing stable for 5+ years; density = moat"],
    ]
    add_mck_table(slide, data, top=Inches(1.7), font_size=Pt(10),
                  col_widths=[Inches(2.5), Inches(1.2), Inches(1.3), Inches(7.3)])


def slide_12_ask(prs):
    """RESOLUTION — The ask + why partner with me."""
    slide = mck_slide(prs,
        "We are seeking $100K in total capital from a co-investor to acquire and operate "
        "a SF-area portable sanitation business with projected 10x+ return over 5 years",
        "Source: Financial model; personal background",
        slide_num=12)

    bullets = [
        ("THE ASK", True),
        "\u2022  Total capital sought: $100K (equity contribution + working capital reserve)",
        "\u2022  Investor contributes $50-75K equity alongside operator's $25-50K",
        "\u2022  Equity split negotiable based on contribution and active/passive role",
        "\u2022  Returns via quarterly distributions + terminal equity at exit (projected 10x+ multiple)",
        "",
        ("TIMELINE", True),
        "\u2022  Mar-Apr 2026: Target identification and LOI    \u2022  May-Jun: Diligence + SBA application",
        "\u2022  Jul-Aug: Close and transition                           \u2022  Sep 2026: Full independent operations",
        "",
        ("WHY PARTNER WITH ME", True),
        "\u2022  Tuck MBA (Dartmouth) — finance, strategy, operations training",
        "\u2022  Chief of Staff experience — built operating systems for scaling companies",
        "\u2022  Tech-first mindset: I'll build the booking/scheduling tools no competitor has",
        "\u2022  Willing to drive the truck Day 1 while building the business — hustle + institutional rigor",
    ]
    add_body_bullets(slide, bullets, font_size=Pt(11))

    add_gold_callout(slide,
        "I'm not buying a job — I'm building an asset.",
        top=Inches(5.9))


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════
def main():
    os.makedirs(OUT_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_cover(prs)
    slide_02_exec_summary(prs)
    slide_03_market(prs)
    slide_04_competitive(prs)
    slide_05_why_now(prs)
    slide_06_business_model(prs)
    slide_07_unit_economics(prs)
    slide_08_projections(prs)
    slide_09_deal_structure(prs)
    slide_10_operating_strategy(prs)
    slide_11_risks(prs)
    slide_12_ask(prs)

    prs.save(OUT_FILE)
    print(f"OK: {OUT_FILE} ({os.path.getsize(OUT_FILE):,} bytes)")


if __name__ == "__main__":
    main()
